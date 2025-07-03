#!/usr/bin/env python3
"""
PowerPoint Template Generator for Sustainability Lab
Creates a professional presentation template with consistent branding
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_lab_template():
    """Create PowerPoint template with Sustainability Lab branding"""
    
    # Create new presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)
    
    # Define colors
    lab_black = RGBColor(0, 0, 0)
    lab_gray = RGBColor(64, 64, 64)
    lab_light_gray = RGBColor(128, 128, 128)
    white = RGBColor(255, 255, 255)
    
    # Slide 1: Title Slide
    title_layout = prs.slide_layouts[0]  # Title slide layout
    slide1 = prs.slides.add_slide(title_layout)
    
    title = slide1.shapes.title
    title.text = "Deep Learning for Computer Vision"
    title.text_frame.paragraphs[0].font.name = "Helvetica Neue"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = lab_black
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle = slide1.placeholders[1]
    subtitle.text = "Advanced Neural Network Architectures\n\nNipun Batra\nSustainability Lab\nJuly 2025"
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.name = "Helvetica Neue"
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = lab_gray
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Add logo placeholder (manually add logo later)
    logo_placeholder = slide1.shapes.add_textbox(
        Inches(11.5), Inches(6), Inches(1.5), Inches(1)
    )
    logo_text = logo_placeholder.text_frame
    logo_text.text = "[LOGO]"
    logo_text.paragraphs[0].font.size = Pt(12)
    logo_text.paragraphs[0].font.color.rgb = lab_light_gray
    logo_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 2: Content Slide Template
    content_layout = prs.slide_layouts[1]  # Title and content layout
    slide2 = prs.slides.add_slide(content_layout)
    
    slide2.shapes.title.text = "Research Motivation"
    slide2.shapes.title.text_frame.paragraphs[0].font.name = "Helvetica Neue"
    slide2.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    slide2.shapes.title.text_frame.paragraphs[0].font.color.rgb = lab_black
    
    # Add horizontal rule under title
    rule_shape = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0.5), Inches(1.8), 
        Inches(12.33), Pt(2)
    )
    rule_shape.fill.solid()
    rule_shape.fill.fore_color.rgb = lab_light_gray
    rule_shape.line.fill.background()
    
    content = slide2.shapes.placeholders[1]
    content.text = """• Computer vision has transformed AI applications
• Deep learning architectures continue to evolve  
• Performance gains through novel architectural innovations
• Real-world deployment challenges remain significant

Key Research Question:
How can we design efficient neural architectures that maintain high accuracy while reducing computational requirements?"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.name = "Helvetica Neue"
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = lab_black
        paragraph.space_after = Pt(12)
    
    # Add logo placeholder to content slide
    content_logo = slide2.shapes.add_textbox(
        Inches(11.5), Inches(0.2), Inches(1.5), Inches(0.8)
    )
    content_logo_text = content_logo.text_frame
    content_logo_text.text = "[LOGO]"
    content_logo_text.paragraphs[0].font.size = Pt(10)
    content_logo_text.paragraphs[0].font.color.rgb = lab_light_gray
    content_logo_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add slide number placeholder
    slide_num = slide2.shapes.add_textbox(
        Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.5)
    )
    slide_num_text = slide_num.text_frame
    slide_num_text.text = "2 / 12"
    slide_num_text.paragraphs[0].font.name = "Helvetica Neue"
    slide_num_text.paragraphs[0].font.size = Pt(14)
    slide_num_text.paragraphs[0].font.color.rgb = lab_light_gray
    slide_num_text.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    # Slide 3: Two-column layout
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide3 = prs.slides.add_slide(blank_layout)
    
    # Title
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(1))
    title3.text_frame.text = "Experimental Setup"
    title3.text_frame.paragraphs[0].font.name = "Helvetica Neue"
    title3.text_frame.paragraphs[0].font.size = Pt(36)
    title3.text_frame.paragraphs[0].font.color.rgb = lab_black
    
    # Horizontal rule
    rule3 = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0.5), Inches(1.3), 
        Inches(12.33), Pt(2)
    )
    rule3.fill.solid()
    rule3.fill.fore_color.rgb = lab_light_gray
    rule3.line.fill.background()
    
    # Left column
    left_col = slide3.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6), Inches(4.5))
    left_col.text_frame.text = """Datasets Used:
• ImageNet-1K (1.28M images)
• CIFAR-10/100
• Custom industrial dataset

Hardware:
• 8× NVIDIA A100 GPUs
• 512GB RAM
• NVMe SSD storage"""
    
    for paragraph in left_col.text_frame.paragraphs:
        paragraph.font.name = "Helvetica Neue"
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = lab_black
        paragraph.space_after = Pt(10)
    
    # Right column with figure placeholder
    right_col = slide3.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5.5), Inches(4.5))
    right_col.text_frame.text = """[NETWORK ARCHITECTURE DIAGRAM]

Input Layer
    ↓
Neural Network Architecture  
    ↓
Output Layer

Figure: Network Overview"""
    
    for paragraph in right_col.text_frame.paragraphs:
        paragraph.font.name = "Helvetica Neue"
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = lab_gray
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Add logo and slide number
    logo3 = slide3.shapes.add_textbox(Inches(11.5), Inches(0.2), Inches(1.5), Inches(0.8))
    logo3.text_frame.text = "[LOGO]"
    logo3.text_frame.paragraphs[0].font.size = Pt(10)
    logo3.text_frame.paragraphs[0].font.color.rgb = lab_light_gray
    logo3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    slide_num3 = slide3.shapes.add_textbox(Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.5))
    slide_num3.text_frame.text = "3 / 12"
    slide_num3.text_frame.paragraphs[0].font.name = "Helvetica Neue"
    slide_num3.text_frame.paragraphs[0].font.size = Pt(14)
    slide_num3.text_frame.paragraphs[0].font.color.rgb = lab_light_gray
    slide_num3.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    # Slide 4: Table slide
    slide4 = prs.slides.add_slide(blank_layout)
    
    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(1))
    title4.text_frame.text = "Performance Comparison"
    title4.text_frame.paragraphs[0].font.name = "Helvetica Neue"
    title4.text_frame.paragraphs[0].font.size = Pt(36)
    title4.text_frame.paragraphs[0].font.color.rgb = lab_black
    
    # Horizontal rule
    rule4 = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0.5), Inches(1.3), 
        Inches(12.33), Pt(2)
    )
    rule4.fill.solid()
    rule4.fill.fore_color.rgb = lab_light_gray
    rule4.line.fill.background()
    
    # Table placeholder
    table_box = slide4.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.5), Inches(3))
    table_box.text_frame.text = """[TABLE: Performance Comparison]

Model               ImageNet Top-1    FLOPs (G)    Parameters (M)
ResNet-50          76.15%            4.1          25.6
EfficientNet-B0    77.32%            0.39         5.3
Our Method         78.94%            0.31         4.2
Vision Transformer 81.28%            17.6         86.4

Key Results:
• Our approach achieves 2.8× fewer FLOPs than ResNet-50
• Maintains competitive accuracy with modern architectures  
• Significant reduction in parameter count enables mobile deployment"""
    
    for paragraph in table_box.text_frame.paragraphs:
        paragraph.font.name = "Helvetica Neue"
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = lab_black
        paragraph.space_after = Pt(8)
    
    # Add logo and slide number
    logo4 = slide4.shapes.add_textbox(Inches(11.5), Inches(0.2), Inches(1.5), Inches(0.8))
    logo4.text_frame.text = "[LOGO]"
    logo4.text_frame.paragraphs[0].font.size = Pt(10)
    logo4.text_frame.paragraphs[0].font.color.rgb = lab_light_gray
    logo4.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    slide_num4 = slide4.shapes.add_textbox(Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.5))
    slide_num4.text_frame.text = "4 / 12"
    slide_num4.text_frame.paragraphs[0].font.name = "Helvetica Neue"
    slide_num4.text_frame.paragraphs[0].font.size = Pt(14)
    slide_num4.text_frame.paragraphs[0].font.color.rgb = lab_light_gray
    slide_num4.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    # Slide 5: Thank You slide
    slide5 = prs.slides.add_slide(title_layout)
    
    thank_title = slide5.shapes.title
    thank_title.text = "Thank You!"
    thank_title.text_frame.paragraphs[0].font.name = "Helvetica Neue"
    thank_title.text_frame.paragraphs[0].font.size = Pt(48)
    thank_title.text_frame.paragraphs[0].font.color.rgb = lab_black
    thank_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    thank_content = slide5.placeholders[1]
    thank_content.text = """Questions & Discussion

Contact: nipun.batra@iiitd.ac.in
Lab Website: https://sustainabilitylab.org
Code: https://github.com/sustainability-lab

[LOGO]"""
    
    for paragraph in thank_content.text_frame.paragraphs:
        paragraph.font.name = "Helvetica Neue"
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = lab_gray
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.space_after = Pt(15)
    
    return prs

def main():
    """Generate the PowerPoint template"""
    print("Creating Sustainability Lab PowerPoint template...")
    
    # Create the template
    prs = create_lab_template()
    
    # Save the template
    output_path = "sustainability_lab_template.pptx"
    prs.save(output_path)
    
    print(f"✅ Template created successfully: {output_path}")
    print("\nTo use this template:")
    print("1. Open the .pptx file in PowerPoint")
    print("2. Replace [LOGO] placeholders with your actual lab logo")
    print("3. Update slide numbers as needed")
    print("4. Save as a PowerPoint template (.potx) for reuse")
    print("\nTemplate features:")
    print("• 16:9 widescreen format")
    print("• Consistent Helvetica fonts")
    print("• Professional gray color scheme")
    print("• Sustainability Lab branding")
    print("• Multiple slide layouts (title, content, two-column, table)")

if __name__ == "__main__":
    main()